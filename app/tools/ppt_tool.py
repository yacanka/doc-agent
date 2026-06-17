"""Safe PPTX inspection and text replacement utilities."""
from __future__ import annotations

import shutil
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation


@dataclass(frozen=True)
class SlideSummary:
    """Compact description of a PowerPoint slide."""

    index: int
    layout_name: str
    shape_count: int
    texts: tuple[str, ...]
    notes: str | None


def inspect_presentation(path: Path) -> tuple[SlideSummary, ...]:
    """Return slide layout, shape, text, and speaker-note information."""
    presentation = Presentation(path)
    return tuple(_summarize_slide(index, slide) for index, slide in enumerate(presentation.slides, start=1))


def extract_text(path: Path) -> str:
    """Extract text from slide shapes and speaker notes where supported."""
    summaries = inspect_presentation(path)
    parts = [text for summary in summaries for text in summary.texts]
    parts.extend(summary.notes or "" for summary in summaries)
    return "\n".join(part for part in parts if part)


def replace_text(source: Path, destination: Path, replacements: dict[str, str]) -> int:
    """Copy a PPTX and replace text inside existing runs when possible."""
    shutil.copy2(source, destination)
    presentation = Presentation(destination)
    changed_count = _replace_presentation_text(presentation, replacements)
    presentation.save(destination)
    return changed_count


def _summarize_slide(index: int, slide: object) -> SlideSummary:
    texts = tuple(_shape_texts(slide.shapes))
    return SlideSummary(index, slide.slide_layout.name, len(slide.shapes), texts, _notes_text(slide))


def _replace_presentation_text(presentation: Presentation, replacements: dict[str, str]) -> int:
    changed_count = 0
    for slide in presentation.slides:
        changed_count += _replace_shape_collection(slide.shapes, replacements)
        changed_count += _replace_notes(slide, replacements)
    return changed_count


def _replace_shape_collection(shapes: object, replacements: dict[str, str]) -> int:
    changed_count = 0
    for shape in shapes:
        changed_count += _replace_shape(shape, replacements)
    return changed_count


def _replace_shape(shape: object, replacements: dict[str, str]) -> int:
    if hasattr(shape, "shapes"):
        return _replace_shape_collection(shape.shapes, replacements)
    if not getattr(shape, "has_text_frame", False):
        return 0
    return _replace_text_frame(shape.text_frame, replacements)


def _replace_notes(slide: object, replacements: dict[str, str]) -> int:
    if not getattr(slide, "has_notes_slide", False):
        return 0
    return _replace_text_frame(slide.notes_slide.notes_text_frame, replacements)


def _replace_text_frame(text_frame: object, replacements: dict[str, str]) -> int:
    return sum(_replace_paragraph(paragraph, replacements) for paragraph in text_frame.paragraphs)


def _replace_paragraph(paragraph: object, replacements: dict[str, str]) -> int:
    return sum(_replace_run(run, replacements) for run in paragraph.runs)


def _replace_run(run: object, replacements: dict[str, str]) -> int:
    changed_count = 0
    for target, replacement in replacements.items():
        occurrences = run.text.count(target)
        if occurrences:
            run.text = run.text.replace(target, replacement)
            changed_count += occurrences
    return changed_count


def _shape_texts(shapes: object) -> list[str]:
    texts: list[str] = []
    for shape in shapes:
        texts.extend(_shape_text(shape))
    return texts


def _shape_text(shape: object) -> list[str]:
    if hasattr(shape, "shapes"):
        return _shape_texts(shape.shapes)
    return [shape.text] if getattr(shape, "has_text_frame", False) and shape.text else []


def _notes_text(slide: object) -> str | None:
    if not getattr(slide, "has_notes_slide", False):
        return None
    text = slide.notes_slide.notes_text_frame.text
    return text or None
