from pathlib import Path

import pytest

fitz = pytest.importorskip("fitz")
openpyxl = pytest.importorskip("openpyxl")
pptx = pytest.importorskip("pptx")

from openpyxl import Workbook, load_workbook
from pptx import Presentation

from app.tools import excel_tool, pdf_tool, ppt_tool


def _xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet["A1"] = "Name"
    sheet["B1"] = "Total"
    sheet["A2"] = "Alice"
    sheet["B2"] = "=SUM(1,2)"
    sheet.merge_cells("C1:D1")
    sheet.column_dimensions["A"].width = 22
    sheet.row_dimensions[1].height = 30
    sheet.auto_filter.ref = "A1:B2"
    sheet.freeze_panes = "A2"
    workbook.properties.creator = "Tester"
    workbook.save(path)


def _pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Hello Alice"
    run = slide.placeholders[1].text_frame.paragraphs[0].add_run()
    run.text = "Project Alice"
    presentation.save(path)


def _pdf(path: Path) -> None:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Read only Alice")
    document.save(path)
    document.close()


def test_excel_inspect_update_append_and_copy_style(tmp_path: Path) -> None:
    source = tmp_path / "source.xlsx"
    updated = tmp_path / "updated.xlsx"
    appended = tmp_path / "appended.xlsx"
    styled = tmp_path / "styled.xlsx"
    _xlsx(source)

    summary = excel_tool.inspect_workbook(source)
    changed = excel_tool.update_cells(source, updated, {"Data": {"A2": "Bob"}})
    appended_count = excel_tool.append_rows(updated, appended, "Data", [["Cara", "=SUM(2,3)"]])
    excel_tool.copy_cell_style(appended, styled, "Data", "A1", "A3")

    workbook = load_workbook(styled, data_only=False)
    sheet = workbook["Data"]
    assert summary.creator == "Tester"
    assert summary.sheets[0].merged_ranges == ("C1:D1",)
    assert changed == 1
    assert appended_count == 1
    assert sheet["B2"].value == "=SUM(1,2)"
    assert sheet["B3"].value == "=SUM(2,3)"
    assert sheet.freeze_panes == "A2"
    assert sheet.auto_filter.ref == "A1:B2"
    assert sheet.column_dimensions["A"].width == 22
    assert sheet.row_dimensions[1].height == 30


def test_ppt_inspect_and_replace_preserves_existing_file(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    destination = tmp_path / "destination.pptx"
    _pptx(source)

    summaries = ppt_tool.inspect_presentation(source)
    changed = ppt_tool.replace_text(source, destination, {"Alice": "Bob"})

    assert summaries[0].layout_name
    assert changed == 2
    assert "Hello Bob" in ppt_tool.extract_text(destination)
    assert "Project Bob" in ppt_tool.extract_text(destination)
    assert "Alice" in ppt_tool.extract_text(source)


def test_pdf_extracts_text_and_rejects_writes(tmp_path: Path) -> None:
    source = tmp_path / "source.pdf"
    _pdf(source)

    assert "Read only Alice" in pdf_tool.extract_text(source)
    with pytest.raises(NotImplementedError):
        pdf_tool.update_pdf(source, {"Alice": "Bob"})
