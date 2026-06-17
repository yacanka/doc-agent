from __future__ import annotations

import hashlib
import json
import zipfile
from pathlib import Path

import pytest

fitz = pytest.importorskip("fitz")
openpyxl = pytest.importorskip("openpyxl")
pptx = pytest.importorskip("pptx")
pytest.importorskip("docx")

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from app.tools import excel_tool, pdf_tool, ppt_tool, word_tool
from app.tools.validation_tool import validate_output_file

from tests.fixtures.document_samples import build_document_fixtures


@pytest.fixture()
def fixtures_directory(tmp_path: Path) -> Path:
    """Create deterministic document samples for one test."""
    return build_document_fixtures(tmp_path / "fixtures")


def _hash(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _zip_entries(path: Path) -> set[str]:
    with zipfile.ZipFile(path) as package:
        return set(package.namelist())


def _changed_entries(before: Path, after: Path) -> set[str]:
    with zipfile.ZipFile(before) as before_zip, zipfile.ZipFile(after) as after_zip:
        names = set(before_zip.namelist()) & set(after_zip.namelist())
        return {name for name in names if before_zip.read(name) != after_zip.read(name)}


@pytest.mark.parametrize("filename", ["sample.docx", "sample.xlsx", "sample.pptx", "sample.pdf"])
def test_original_fixture_hashes_remain_unchanged(tmp_path: Path, fixtures_directory: Path, filename: str) -> None:
    original = fixtures_directory / filename
    before_hash = _hash(original)
    working_copy = tmp_path / filename
    working_copy.write_bytes(original.read_bytes())

    if filename.endswith(".docx"):
        word_tool.replace_text(original, tmp_path / "edited.docx", {"Alice": "Bob"})
    elif filename.endswith(".xlsx"):
        excel_tool.update_cells(original, tmp_path / "edited.xlsx", {"Data": {"A2": "Bob"}})
    elif filename.endswith(".pptx"):
        ppt_tool.replace_text(original, tmp_path / "edited.pptx", {"Alice": "Bob"})
    else:
        pdf_tool.extract_text(original)

    assert _hash(original) == before_hash
    assert _hash(working_copy) == before_hash


def test_docx_text_extraction_and_single_run_style_preservation(tmp_path: Path, fixtures_directory: Path) -> None:
    source = fixtures_directory / "sample.docx"
    destination = tmp_path / "single.docx"

    changed = word_tool.replace_text(source, destination, {"Hello Alice": "Hello Bob"})
    document = Document(destination)
    run = document.paragraphs[0].runs[0]

    assert "Extract me from DOCX" in word_tool.extract_text(source)
    assert changed == 1
    assert run.text == "Hello Bob"
    assert run.bold is True
    assert run.italic is True
    assert _zip_entries(source) == _zip_entries(destination)
    assert _changed_entries(source, destination) <= {"word/document.xml"}


def test_docx_multi_run_replacement_preserves_package_structure(tmp_path: Path, fixtures_directory: Path) -> None:
    source = fixtures_directory / "sample.docx"
    requested = tmp_path / "multi.docx"

    result = word_tool.replace_text_preserve_style(source, "Alice", "Bob", requested)
    document = Document(result.output_path)
    second_paragraph = document.paragraphs[1]

    assert result.changed_count == 1
    assert second_paragraph.text == "Multi Bob"
    assert len(second_paragraph.runs) == 3
    assert _zip_entries(source) == _zip_entries(result.output_path)
    assert _changed_entries(source, result.output_path) <= {"word/document.xml"}


def test_xlsx_cell_update_preserves_styles_formulas_and_structure(tmp_path: Path, fixtures_directory: Path) -> None:
    source = fixtures_directory / "sample.xlsx"
    destination = tmp_path / "updated.xlsx"

    changed = excel_tool.update_cells(source, destination, {"Data": {"A2": "Bob"}})
    workbook = load_workbook(destination, data_only=False)
    sheet = workbook["Data"]

    assert changed == 1
    assert sheet["A2"].value == "Bob"
    assert sheet["A2"]._style == sheet["B2"]._style
    assert sheet["C2"].value == "=SUM(B2:B2)"
    assert sheet.freeze_panes == "A2"
    assert sheet.auto_filter.ref == "A1:C2"
    assert _zip_entries(source) == _zip_entries(destination)
    assert not (_changed_entries(source, destination) - {"xl/worksheets/sheet1.xml", "docProps/core.xml"})


def test_xlsx_append_row_copies_style(tmp_path: Path, fixtures_directory: Path) -> None:
    source = fixtures_directory / "sample.xlsx"
    destination = tmp_path / "appended.xlsx"

    count = excel_tool.append_rows(source, destination, "Data", [["Cara", 20, "=SUM(B3:B3)"]])
    sheet = load_workbook(destination, data_only=False)["Data"]

    assert count == 1
    assert sheet["A3"].value == "Cara"
    assert sheet["C3"].value == "=SUM(B3:B3)"
    assert sheet["A3"]._style == sheet["A2"]._style
    assert sheet.row_dimensions[3].height == sheet.row_dimensions[2].height
    assert _zip_entries(source) == _zip_entries(destination)


def test_pptx_text_replacement_preserves_shape_and_font_properties(tmp_path: Path, fixtures_directory: Path) -> None:
    source = fixtures_directory / "sample.pptx"
    destination = tmp_path / "edited.pptx"
    before = Presentation(source).slides[0].shapes[0]

    changed = ppt_tool.replace_text(source, destination, {"Alice": "Bob"})
    after = Presentation(destination).slides[0].shapes[0]
    run = after.text_frame.paragraphs[0].runs[0]

    assert changed == 1
    assert "Hello Bob" in ppt_tool.extract_text(destination)
    assert (after.left, after.top, after.width, after.height) == (before.left, before.top, before.width, before.height)
    assert run.font.name == "Aptos"
    assert run.font.size.pt == 28
    assert run.font.bold is True
    assert _zip_entries(source) == _zip_entries(destination)
    assert _changed_entries(source, destination) <= {"ppt/slides/slide1.xml"}


def test_pdf_read_only_extraction_does_not_mutate_fixture(fixtures_directory: Path) -> None:
    source = fixtures_directory / "sample.pdf"
    before_hash = _hash(source)

    assert "Read only PDF Alice" in pdf_tool.extract_text(source)
    with pytest.raises(NotImplementedError):
        pdf_tool.update_pdf(source, {"Alice": "Bob"})
    assert _hash(source) == before_hash


def test_validation_report_generation_and_timestamped_output(tmp_path: Path, fixtures_directory: Path) -> None:
    source = fixtures_directory / "sample.docx"
    requested = tmp_path / "result.docx"

    result = word_tool.replace_text_preserve_style(source, "Alice", "Bob", requested)
    validation = validate_output_file(source, result.output_path, tmp_path / "workspace")
    report = json.loads(Path(validation.report_path).read_text(encoding="utf-8"))

    assert result.output_path != requested
    assert result.output_path.name.endswith("_result.docx")
    assert validation.valid is True
    assert validation.parser == "python-docx"
    assert report["created_at"]
    assert report["changed_entries"] == ["word/document.xml"]
