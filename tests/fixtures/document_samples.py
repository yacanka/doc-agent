"""Deterministic document fixture builders for regression tests."""
from __future__ import annotations

from pathlib import Path


def build_document_fixtures(directory: Path) -> Path:
    """Create small DOCX, XLSX, PPTX, and PDF samples inside the directory."""
    directory.mkdir(parents=True, exist_ok=True)
    _build_docx(directory / "sample.docx")
    _build_xlsx(directory / "sample.xlsx")
    _build_pptx(directory / "sample.pptx")
    _build_pdf(directory / "sample.pdf")
    return directory


def _build_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    run = document.add_paragraph().add_run("Hello Alice")
    run.bold = True
    run.italic = True
    paragraph = document.add_paragraph()
    paragraph.add_run("Multi ")
    paragraph.add_run("Al").bold = True
    paragraph.add_run("ice").italic = True
    document.add_paragraph("Extract me from DOCX")
    document.save(path)


def _build_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    _populate_xlsx_values(sheet)
    _style_xlsx_header(sheet)
    _style_xlsx_data_row(sheet)
    sheet.freeze_panes = "A2"
    sheet.auto_filter.ref = "A1:C2"
    sheet.column_dimensions["A"].width = 20
    workbook.save(path)


def _populate_xlsx_values(sheet: object) -> None:
    sheet["A1"] = "Name"
    sheet["B1"] = "Value"
    sheet["C1"] = "Total"
    sheet["A2"] = "Alice"
    sheet["B2"] = 10
    sheet["C2"] = "=SUM(B2:B2)"


def _style_xlsx_header(sheet: object) -> None:
    from openpyxl.styles import Alignment, Font, PatternFill

    for cell in sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="1F4E79")
        cell.alignment = Alignment(horizontal="center")


def _style_xlsx_data_row(sheet: object) -> None:
    from openpyxl.styles import Border, Font, PatternFill, Side

    border = Border(bottom=Side(style="thin"))
    for cell in sheet[2]:
        cell.font = Font(name="Calibri", size=12)
        cell.fill = PatternFill("solid", fgColor="E2F0D9")
        cell.border = border
    sheet.row_dimensions[2].height = 24


def _build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Hello Alice"
    run.font.name = "Aptos"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    presentation.save(path)


def _build_pdf(path: Path) -> None:
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Read only PDF Alice")
    document.save(path)
    document.close()
